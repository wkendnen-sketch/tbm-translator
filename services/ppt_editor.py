import os
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt


MAX_FONT_SIZES = {
    "ko": 22,
    "zh": 22,
    "vi": 20,
    "my": 18,
}


def find_photo_box(slide):
    """
    이름이 PHOTO_BOX 인 도형 찾기
    """
    for shape in slide.shapes:
        try:
            if shape.name == "PHOTO_BOX":
                return shape
        except Exception:
            pass

    raise ValueError("PHOTO_BOX 도형을 찾지 못했습니다.")


def remove_shape(shape):
    sp = shape._element
    sp.getparent().remove(sp)


def insert_photo(slide, image_path, box):
    """
    PHOTO_BOX 크기 그대로 사진 삽입
    """
    slide.shapes.add_picture(
        image_path,
        box.left,
        box.top,
        width=box.width,
        height=box.height,
    )


def _calc_font_size(text, max_size):
    text_len = len(str(text).strip())

    if text_len <= 8:
        font_size = max_size
    elif text_len <= 14:
        font_size = max_size - 2
    elif text_len <= 22:
        font_size = max_size - 4
    elif text_len <= 32:
        font_size = max_size - 6
    else:
        font_size = max_size - 8

    return max(font_size, 10)


def set_cell_text_and_fit(cell, text, max_size):
    text = str(text).strip()

    cell.text = text

    tf = cell.text_frame
    tf.word_wrap = True

    font_size = _calc_font_size(text, max_size)

    for paragraph in tf.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(font_size)


def find_table(slide):
    for shape in slide.shapes:
        if hasattr(shape, "table"):
            return shape.table
    raise ValueError("표를 찾지 못했습니다.")


def fill_template_ppt(template_path, image_paths, translated_items):

    prs = Presentation(template_path)

    if len(image_paths) > len(prs.slides):
        raise ValueError("사진 수가 슬라이드 수보다 많습니다.")

    for idx, (img_path, item) in enumerate(zip(image_paths, translated_items)):

        slide = prs.slides[idx]

        photo_box = find_photo_box(slide)

        insert_photo(slide, img_path, photo_box)

        remove_shape(photo_box)

        table = find_table(slide)

        set_cell_text_and_fit(table.cell(2, 1), item["ko"], MAX_FONT_SIZES["ko"])
        set_cell_text_and_fit(table.cell(3, 1), item["zh"], MAX_FONT_SIZES["zh"])
        set_cell_text_and_fit(table.cell(4, 1), item["vi"], MAX_FONT_SIZES["vi"])
        set_cell_text_and_fit(table.cell(5, 1), item["my"], MAX_FONT_SIZES["my"])

    output_dir = "output"
    os.makedirs(output_dir, exist_ok=True)

    output_path = Path(output_dir) / f"result_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"

    prs.save(str(output_path))

    return str(output_path)